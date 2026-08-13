import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9 (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    DARK_BG = RGBColor(8, 14, 11)       # #080E0B
    CARD_BG = RGBColor(18, 28, 23)     # #121C17
    LIME_ACCENT = RGBColor(212, 255, 51) # #D4FF33
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(180, 190, 185)
    
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_header(slide, slide_num, title_text, category_text="SUNVEXA CLEAN-TECH PLATFORM"):
        # Category / Slide Num
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = cat_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{category_text}  •  SLIDE {slide_num:02d}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = LIME_ACCENT
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    blank_layout = prs.slide_layouts[6]

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_dark_background(s1)
    
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "SUNVEXA"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = LIME_ACCENT
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Smarter Solar. Brighter Future."
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    p3 = tf.add_paragraph()
    p3.text = "\nNext-Generation Clean-Tech E-Commerce & Solar Systems Platform"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "\nReact 18  •  HTML5 240-Frame Canvas  •  Java 21  •  Spring Boot 3  •  PostgreSQL 18"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = LIME_ACCENT
    p4.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 2: Executive Summary & Project Vision
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_dark_background(s2)
    add_header(s2, 2, "Executive Summary & Project Vision")

    # Card 1: Problem
    c1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    c1.fill.solid(); c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = RGBColor(239, 68, 68)
    tf1 = c1.text_frame; tf1.word_wrap = True
    p = tf1.paragraphs[0]; p.text = "⚠️ THE INDUSTRY CHALLENGE"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = RGBColor(239, 68, 68)
    items1 = [
        "Traditional solar websites are static, plain, and lack interactive price transparency.",
        "Rooftop capacity & ROI calculation is mathematically complex for general buyers.",
        "Lack of real-time installation scheduling and order fulfillment tracking.",
        "Unsecured checkout processes without verifiable database persistence."
    ]
    for item in items1:
        p = tf1.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(13); p.font.color.rgb = WHITE

    # Card 2: Solution
    c2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    c2.fill.solid(); c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = LIME_ACCENT
    tf2 = c2.text_frame; tf2.word_wrap = True
    p = tf2.paragraphs[0]; p.text = "✨ THE SUNVEXA PLATFORM"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = LIME_ACCENT
    items2 = [
        "Immersive 240-frame 3D solar rendering canvas executing at 60 FPS.",
        "AI Solar Copilot & 24H Load Simulator for automated capacity & ROI estimation.",
        "Transparent 6-Step Buy Now flow with turnkey installation options (+₹25,000).",
        "Enterprise Java 21 / Spring Boot 3 REST APIs backed by PostgreSQL 18 DB."
    ]
    for item in items2:
        p = tf2.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(13); p.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 3: 5 Key Clean-Tech Innovations
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_dark_background(s3)
    add_header(s3, 3, "5 Core Clean-Tech Innovations")

    features = [
        ("🎬 240-Frame Canvas Scroll Animation", "Smooth 60 FPS 3D solar array sequence synchronized seamlessly to scroll position."),
        ("🤖 AI Solar Copilot Chat Engine", "Conversational AI recommending panel wattages, hybrid inverters & 25-year financial savings."),
        ("🛰️ RoofScan AI Satellite Analysis", "Preliminary rooftop surface area, tilt angle & 92% sun irradiance index assessment."),
        ("⚡ 24-Hour Energy Simulator", "Hourly load profiling comparing solar generation vs battery storage & grid net-metering."),
        ("🛠️ Custom Solar System Builder", "Drag-and-drop hardware assembly (Panels, Inverters, Batteries) with real-time specs checking.")
    ]

    for i, (ft_title, ft_desc) in enumerate(features):
        row = i // 3
        col = i % 3
        x = Inches(0.8 + col * 3.9)
        y = Inches(1.8 + row * 2.5)
        shape = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.2))
        shape.fill.solid(); shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = LIME_ACCENT
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = ft_title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = LIME_ACCENT
        p2 = tf.add_paragraph(); p2.text = f"\n{ft_desc}"; p2.font.size = Pt(11); p2.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 4: Full-Stack System Architecture
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_dark_background(s4)
    add_header(s4, 4, "Full-Stack System Architecture")

    box_w = Inches(3.6); box_h = Inches(4.8)
    
    # Layer 1: Frontend
    lay1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), box_w, box_h)
    lay1.fill.solid(); lay1.fill.fore_color.rgb = CARD_BG; lay1.line.color.rgb = LIME_ACCENT
    tf1 = lay1.text_frame; tf1.word_wrap = True
    p = tf1.paragraphs[0]; p.text = "FRONTEND LAYER\nReact 18 + TypeScript"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = LIME_ACCENT
    f_items = ["Vite Build Engine", "HTML5 Canvas API", "Tailwind CSS Design System", "Lucide Icon Suite", "Resilient API Client (src/services/api.ts)"]
    for fit in f_items:
        p = tf1.add_paragraph(); p.text = f"• {fit}"; p.font.size = Pt(12); p.font.color.rgb = WHITE

    # Layer 2: Backend
    lay2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.85), Inches(1.8), box_w, box_h)
    lay2.fill.solid(); lay2.fill.fore_color.rgb = CARD_BG; lay2.line.color.rgb = LIME_ACCENT
    tf2 = lay2.text_frame; tf2.word_wrap = True
    p = tf2.paragraphs[0]; p.text = "BACKEND LAYER\nJava 21 + Spring Boot 3"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = LIME_ACCENT
    b_items = ["Spring Web REST Controllers", "Spring Security + JJWT Auth", "BCrypt Password Encoder", "Spring Data JPA / Hibernate", "Global Exception Handler (@RestControllerAdvice)"]
    for bit in b_items:
        p = tf2.add_paragraph(); p.text = f"• {bit}"; p.font.size = Pt(12); p.font.color.rgb = WHITE

    # Layer 3: Database
    lay3 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.8), box_w, box_h)
    lay3.fill.solid(); lay3.fill.fore_color.rgb = CARD_BG; lay3.line.color.rgb = LIME_ACCENT
    tf3 = lay3.text_frame; tf3.word_wrap = True
    p = tf3.paragraphs[0]; p.text = "DATABASE LAYER\nPostgreSQL 18"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = LIME_ACCENT
    d_items = ["Flyway DDL Schema Versioning", "11 Relational Tables (sunvexa_db)", "Indexed FK References", "Stateful Price Snapshotting", "Seed Catalog Data (V2__seed_data.sql)"]
    for dit in d_items:
        p = tf3.add_paragraph(); p.text = f"• {dit}"; p.font.size = Pt(12); p.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 5: 6-Step Purchase Flow
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_dark_background(s5)
    add_header(s5, 5, "Transparent 6-Step Purchase Engine")

    steps = [
        ("Step 1: Product Review", "Review solar specs, power output, warranty terms, and unit price."),
        ("Step 2: Customer Details", "Collect full name, email, phone, shipping address, city, state & pin code."),
        ("Step 3: Delivery & Install", "Select Product Only OR Turnkey Rooftop Site Installation (+₹25,000)."),
        ("Step 4: Demo Payment", "Process payment via UPI, Credit/Debit Card, or Net Banking with reference ID."),
        ("Step 5: Order Confirmation", "Persists order record in PostgreSQL sunvexa_db with stateful price snapshots."),
        ("Step 6: Live Order Tracking", "Real-time 6-phase timeline tracking progress by order number (e.g. SNR-89421).")
    ]

    for i, (st_title, st_desc) in enumerate(steps):
        row = i // 3
        col = i % 3
        x = Inches(0.8 + col * 3.9)
        y = Inches(1.8 + row * 2.5)
        shape = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.2))
        shape.fill.solid(); shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = LIME_ACCENT
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = st_title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = LIME_ACCENT
        p2 = tf.add_paragraph(); p2.text = f"\n{st_desc}"; p2.font.size = Pt(11); p2.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 6: Security & Authentication Engine
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_dark_background(s6)
    add_header(s6, 6, "Security, JWT & Database Authentication")

    sec_items = [
        ("🔐 JWT Bearer Token Security", "Stateless Authentication using HMAC-SHA512 signed JWT tokens passed via Authorization: Bearer header."),
        ("🛡️ BCrypt Password Encryption", "Passwords encrypted using BCrypt strength 10 prior to persistence in PostgreSQL users table."),
        ("👥 Role-Based Access Control (RBAC)", "Distinguishes CUSTOMER and ADMIN roles using Spring Security @PreAuthorize('hasRole(\"ADMIN\")')."),
        ("🔍 Explicit Verification Logic", "Strict backend checks for duplicate emails on registration and invalid password detection on login.")
    ]

    for i, (title, desc) in enumerate(sec_items):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 2.5)
        shape = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.2))
        shape.fill.solid(); shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = LIME_ACCENT
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = LIME_ACCENT
        p2 = tf.add_paragraph(); p2.text = f"\n{desc}"; p2.font.size = Pt(12); p2.font.color.rgb = WHITE

    # Save presentation
    output_path = os.path.join(os.getcwd(), "SUNVEXA_Full_Presentation.pptx")
    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at: {output_path}")

if __name__ == "__main__":
    create_presentation()
